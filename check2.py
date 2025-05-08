import os

# Set this before importing transformers to avoid TensorFlow conflicts
os.environ["TRANSFORMERS_NO_TF"] = "1"
os.environ["TRANSFORMERS_NO_CODECARBON"] = "1"

import streamlit as st
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
import cv2
import pytesseract
import tempfile
import cohere
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
from tensorflow.keras.applications import MobileNetV2
from tensorflow.keras.applications.mobilenet_v2 import preprocess_input, decode_predictions

from datasets import load_metric
metric = load_metric("accuracy")
metric = load_metric(custom_metric)

def custom_metric(predictions, references):
    # Your custom metric logic here
    return {"accuracy": some_accuracy_value}

# Register custom metric under a unique name
metric = load_metric(custom_metric)

    




@st.cache_resource
def load_embedding_model():
    return SentenceTransformer('all-MiniLM-L6-v2')

model = load_embedding_model()

def split_text(text, chunk_size=500, overlap=100):
    words = text.split()
    chunks = []
    start = 0
    while start < len(words):
        end = min(start + chunk_size, len(words))
        chunk = " ".join(words[start:end])
        chunks.append(chunk)
        start += chunk_size - overlap
    return chunks

def build_faiss_index(chunks):
    embeddings = model.encode(chunks)
    dimension = embeddings.shape[1]
    index = faiss.IndexFlatL2(dimension)
    index.add(embeddings)
    return index, embeddings

def retrieve_relevant_chunks(query, index, chunks, embeddings, top_k=3):
    query_embedding = model.encode([query])
    _, indices = index.search(query_embedding, top_k)
    return [chunks[i] for i in indices[0]]

def convert_pptx_to_text(file):
    prs = Presentation(file)
    return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])

def convert_pdf_to_text(file):
    doc = fitz.open(stream=file.read(), filetype="pdf")
    return "\n".join([page.get_text("text") for page in doc])

def convert_docx_to_text(file):
    doc = Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

def get_file_text(file):
    ext = os.path.splitext(file.name)[1].lower()
    if ext == ".txt":
        return file.read().decode("utf-8")
    elif ext == ".pptx":
        return convert_pptx_to_text(file)
    elif ext == ".docx":
        return convert_docx_to_text(file)
    elif ext == ".pdf":
        return convert_pdf_to_text(file)
    else:
        st.error("Unsupported file format.")
        return None

def initialize_cohere(api_key):
    return cohere.Client(api_key)

def query_cohere(client, question, context):
    try:
        response = client.chat(
            message=question,
            documents=[{"text": context}],
            temperature=0.7,
        )
        return response.text
    except Exception as e:
        return f"Error: {str(e)}"

def load_and_preprocess_image(path):
    image = cv2.imread(path)
    image = cv2.cvtColor(image, cv2.COLOR_BGR2RGB)
    image = cv2.resize(image, (224, 224))
    return preprocess_input(image)

def predict_image(path):
    model = MobileNetV2(weights='imagenet')
    image = load_and_preprocess_image(path)
    image = np.expand_dims(image, axis=0)
    preds = model.predict(image)
    return [f"{label}: {score*100:.2f}%" for (_, label, score) in decode_predictions(preds, top=3)[0]]

def extract_images_from_pdf(file, output_dir):
    doc = fitz.open(stream=file.read(), filetype="pdf")
    images = []
    for i, page in enumerate(doc):
        for j, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base = doc.extract_image(xref)
            image_path = os.path.join(output_dir, f"page_{i+1}_img{j+1}.{base['ext']}")
            with open(image_path, "wb") as f:
                f.write(base["image"])
            images.append(image_path)
    return images

def extract_images_from_docx(file, output_dir):
    doc = Document(file)
    images = []
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            path = os.path.join(output_dir, rel.target_ref.split("/")[-1])
            with open(path, "wb") as f:
                f.write(rel.target_part.blob)
            images.append(path)
    return images

def extract_images_from_pptx(file, output_dir):
    prs = Presentation(file)
    images = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == 13:
                blob = shape.image.blob
                ext = shape.image.ext
                path = os.path.join(output_dir, f"slide_{i+1}.{ext}")
                with open(path, "wb") as f:
                    f.write(blob)
                images.append(path)
    return images

# UI
st.title("Lansera - Document & Image Analyzer")

option = st.sidebar.selectbox("Choose Feature:", [
    "Analyze Document and Ask Questions",
    "Predict Content of an Image",
    "Analyze Document and Find Relevant Images"
])

COHERE_API_KEY = st.secrets["COHERE_API_KEY"] if "COHERE_API_KEY" in st.secrets else os.getenv("COHERE_API_KEY")
if not COHERE_API_KEY:
    st.warning("Cohere API key not found. Please set it in environment or secrets.")

if option == "Analyze Document and Ask Questions":
    uploaded_file = st.file_uploader("Upload document:", type=["txt", "pdf", "pptx", "docx"])
    if uploaded_file and COHERE_API_KEY:
        text = get_file_text(uploaded_file)
        if text:
            chunks = split_text(text)
            index, embeddings = build_faiss_index(chunks)
            question = st.text_input("Your question:")
            if question:
                relevant_chunks = retrieve_relevant_chunks(question, index, chunks, embeddings)
                context = "\n\n".join(relevant_chunks)
                st.write("### Relevant Chunks:")
                for i, ch in enumerate(relevant_chunks):
                    st.markdown(f"**Chunk {i+1}:** {ch}")
                co = initialize_cohere(COHERE_API_KEY)
                st.write("### *Answer:*")
                st.write(query_cohere(co, question, context))


